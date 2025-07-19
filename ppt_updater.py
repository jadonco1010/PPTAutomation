import zipfile
import xml.etree.ElementTree as ET
import logging
import re
import pandas as pd
from openpyxl import Workbook
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.util import Inches # Important for unit conversion
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR, MSO_FILL
from datetime import datetime
from typing import Dict, List, Tuple, Set, Any

from utils import coordinate_to_tuple, find_dynamic_sheets, get_fiscal_quarter_and_month, iterate_all_shapes

# Module-level dictionary to store metadata about percentage bar shapes
_percentage_shapes_metadata: Dict[str, Dict[str, Dict[str, Any]]] = {}

def _find_sheet_xml_path_for_hidden(z: zipfile.ZipFile, sheet_name: str) -> str:
    """Helper to find the XML path for a given sheet within the Excel zip archive (used for hidden rows/cols)."""
    ns_wb  = {"ns":"http://schemas.openxmlformats.org/spreadsheetml/2006/main"}
    wb_xml = ET.fromstring(z.read("xl/workbook.xml"))
    sh     = wb_xml.find(f".//ns:sheet[@name='{sheet_name}']", ns_wb)
    if sh is None:
        raise KeyError(f"Sheet '{sheet_name}' not found")
    rid = sh.attrib["{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"]

    ns_rel = {"pr":"http://schemas.openxmlformats.org/package/2006/relationships"}
    rels   = ET.fromstring(z.read("xl/_rels/workbook.xml.rels"))
    entry  = rels.find(f".//pr:Relationship[@Id='{rid}']", ns_rel)
    if entry is None:
        raise KeyError(f"Relationship {rid} not in workbook.xml.rels")
    target = entry.attrib["Target"].lstrip("/")
    return target

def _get_hidden_rows_cols_from_xml(file_path: str, sheet_name: str) -> Tuple[Set[int], Set[int]]:
    """Extracts hidden row and column numbers for a specific sheet by parsing its XML."""
    with zipfile.ZipFile(file_path) as z:
        xml_path = _find_sheet_xml_path_for_hidden(z, sheet_name)
        tree     = ET.fromstring(z.read(xml_path))
        ns       = {"ns":"http://schemas.openxmlformats.org/spreadsheetml/2006/main"}

        hidden_rows = {
            int(r.attrib["r"])
            for r in tree.findall(".//ns:row", ns)
            if r.get("hidden") == "1"
        }
        hidden_cols = set()
        for col in tree.findall(".//ns:cols/ns:col", ns):
            if col.get("hidden") == "1":
                mn, mx = int(col.get("min")), int(col.get("max"))
                hidden_cols.update(range(mn, mx+1))
    
    return hidden_rows, hidden_cols

def _extract_openpyxl_block(
    workbook: Workbook, sheet_name: str, start_cell: str, end_cell: str,
    hidden_rows: Set[int], hidden_cols: Set[int]
) -> pd.DataFrame:
    """Extracts a block of data from an openpyxl workbook, respecting hidden rows/columns."""
    ws = workbook[sheet_name]
    r1, c1 = coordinate_to_tuple(start_cell)
    r2, c2 = coordinate_to_tuple(end_cell)

    data = []
    for r in range(r1, r2+1):
        if r in hidden_rows:
            continue
        row_vals = []
        for c in range(c1, c2+1):
            if c in hidden_cols:
                continue
            row_vals.append(ws.cell(row=r, column=c).value)
        data.append(row_vals)

    return pd.DataFrame(data)

def load_tables_from_excel(
    file_path: str,
    sheet_names: List[str],
    table_regions: Dict[str, List[Tuple[str, str]]],
    workbook: Workbook # Pass the loaded workbook here
) -> Dict[str, List[pd.DataFrame]]:
    """
    Loads specified table regions from an Excel file into pandas DataFrames,
    handling hidden rows and columns.
    """
    all_tables: Dict[str, List[pd.DataFrame]] = {}

    for sheet in sheet_names:
        try:
            hidden_rows, hidden_cols = _get_hidden_rows_cols_from_xml(file_path, sheet)
        except KeyError:
            logging.warning(f"Sheet '{sheet}' not found for table loading → skipping")
            continue

        dfs: List[pd.DataFrame] = []
        for start, end in table_regions[sheet]:
            r1, c1 = coordinate_to_tuple(start)
            r2, c2 = coordinate_to_tuple(end)

            try:
                # Use pandas read_excel for initial data load
                df = pd.read_excel(
                    file_path,
                    sheet_name=sheet,
                    header=None,
                    skiprows=r1 - 1,
                    nrows=(r2 - r1 + 1),
                    usecols=f"{get_column_letter(c1)}:{get_column_letter(c2)}",
                    engine="openpyxl",
                    dtype=object
                )
            except Exception as e:
                logging.error(f"Error reading range {start}:{end} in sheet '{sheet}': {e}")
                continue

            # Apply hidden row/column filtering
            drop_r = [r - r1 for r in hidden_rows if r1 <= r <= r2]
            drop_r = [r for r in drop_r if r < len(df)]
            if drop_r:
                df = df.drop(index=drop_r).reset_index(drop=True)

            drop_c = [c - c1 for c in hidden_cols if c1 <= c <= c2]
            drop_c = [c for c in drop_c if c < len(df.columns)]
            if drop_c:
                cols_to_drop = df.columns[drop_c]
                df = df.drop(columns=cols_to_drop)

            # Handle formula values by extracting from openpyxl workbook
            fallback = _extract_openpyxl_block(workbook, sheet, start, end, hidden_rows, hidden_cols)
            fallback = fallback.iloc[: df.shape[0], : df.shape[1]] # Ensure dimensions match after drops
            mask = df.map(lambda x: isinstance(x, str) and x.startswith("="))
            for i, j in zip(*mask.to_numpy().nonzero()):
                df.iat[i, j] = fallback.iat[i, j]
            
            pd.set_option('future.no_silent_downcasting', True)
            df = df.fillna("").infer_objects(copy=False)

            dfs.append(df)

        all_tables[sheet] = dfs

    return all_tables

def _get_date_labels() -> Dict[str,str]:
    """Generates a dictionary of date-related labels for PowerPoint placeholders."""
    now = datetime.now()
    fiscal_year, fiscal_quarter_str, fiscal_month_str, fiscal_month_in_quarter_str = get_fiscal_quarter_and_month(now.date())

    actual_month_name = now.strftime("%B")
    actual_day = now.day
    actual_year = now.year
    actual_date_mm_dd_yy = now.strftime("%m.%d.%y")

    return {
        "QuarterLabel":  fiscal_quarter_str,
        "MonthLabel":    fiscal_month_in_quarter_str,
        "Date":          f"{actual_month_name} {actual_day}, {actual_year}",
        "YearLabel":     str(fiscal_year),
        "Title":         f"{fiscal_month_in_quarter_str} & {fiscal_quarter_str} Fcst",
        "dateLabel":     actual_date_mm_dd_yy,
        "Month":         actual_month_name
    }

def _replace_date_tags(prs: Presentation):
    """Replaces date-related placeholders (e.g., {{Date}}) in the PowerPoint presentation."""
    labels = _get_date_labels()

    for slide in prs.slides:
        for shape in iterate_all_shapes(slide.shapes):
            if not getattr(shape, "has_text_frame", False):
                continue

            for paragraph in shape.text_frame.paragraphs:
                full_text = "".join(run.text for run in paragraph.runs)
                new_text = full_text

                for key, val in labels.items():
                    tag = f"{{{{{key}}}}}"
                    if tag in new_text:
                        new_text = new_text.replace(tag, str(val))

                if new_text != full_text:
                    for run in paragraph.runs:
                        run.text = "" # Clear existing runs
                    if paragraph.runs:
                        paragraph.runs[0].text = new_text # Update first run
                    else:
                        paragraph.add_run().text = new_text # Add new run if none exist

def _format_custom_value(prefix: str, value: Any) -> Tuple[str, float]:
    """
    Formats a numeric value based on a given prefix for PowerPoint display.
    Returns a tuple of (formatted_value, raw_percentage_for_shapes).
    raw_percentage_for_shapes is the value normalized to 0.0-1.0 for shape width calculation.
    """
    raw_percentage_for_shapes = 0.0 # Default

    # Handle string values gracefully
    if isinstance(value, str):
        return value, raw_percentage_for_shapes
    
    # Handle None or NaN values
    if value is None or pd.isna(value):
        return "", raw_percentage_for_shapes

    # Convert to float for numeric values
    value = float(value)
    formatted_value = ""

    # Handle millions formatting (a-z, ab-ah)
    if prefix in {"a", "b", "c", "d", "e", "f", "g", "h", "i", "j", "k", "l", "m", 
                  "n", "o", "p", "q", "r", "s", "t", "u", "v", "w", "x", "y", "z", 
                  "ab", "ac", "ad", "ae", "af", "ag", "ah"}:
        formatted_value = f"{value / 1_000_000:.1f}"
        if value < 0:
            formatted_value = f"({abs(value) / 1_000_000:.1f})"
        return formatted_value, raw_percentage_for_shapes

    # These prefixes get integer percentage formatting
    if prefix in {"aa", "bb", "cc", "dd", "ee", "ff", "gg", "hh", "ii"}:
        raw_percentage_for_shapes = value # Use raw value for shape calculation (e.g., 0.56)
        formatted_value = f"{round(value * 100):,}%" # Formats to integer percentage (e.g., 56%)
        return formatted_value, raw_percentage_for_shapes
    
    # These prefixes get one-decimal percentage formatting
    if prefix in {"AB", "BC", "CD", "DE", "EF", "FG", "HH", "II"}:
        raw_percentage_for_shapes = value # Use raw value for shape calculation (e.g., 0.56)
        formatted_value = f"{value * 100:.1f}%" # Formats to one decimal place (e.g., 56.0%)
        return formatted_value, raw_percentage_for_shapes
    
    if prefix in {"A", "B", "C", "D", "E", "F"}:
        return f"{int(value) // 1_000:,}", raw_percentage_for_shapes
    
    # Default return if no specific prefix match
    return str(value), raw_percentage_for_shapes


def _collect_initial_shape_data(prs: Presentation):
    """
    Collects initial width and left position of shapes intended for percentage bars.
    These shapes should be named with the convention 'bar_XXYY' (e.g., 'bar_aa3').
    """
    global _percentage_shapes_metadata
    _percentage_shapes_metadata = {} # Clear previous data

    logging.info("Collecting initial shape data for percentage bars...")
    found_shapes_count = 0
    for slide_idx, slide in enumerate(prs.slides):
        slide_id = slide.slide_id # Use slide_id for unique identification
        _percentage_shapes_metadata[slide_id] = {}
        for shape in iterate_all_shapes(slide.shapes):
            # Check if shape has a name and matches the 'bar_XXYY' convention
            # Removed: 'not getattr(shape, "has_text_frame", False)' condition
            # This allows rectangles (which have text frames) to be identified as bar shapes.
            if hasattr(shape, 'name'):
                match = re.match(r"^bar_([A-Za-z]{1,2}\d+)$", shape.name)
                if match:
                    # Ensure it's not a table itself (which also has a name)
                    if shape.has_table:
                        logging.warning(f"  Shape '{shape.name}' on slide {slide_idx+1} matches naming convention but is a Table. Skipping.")
                        continue
                    
                    xx_yy_key = match.group(1)
                    _percentage_shapes_metadata[slide_id][shape.name] = {
                        "shape_obj": shape,
                        "original_width": shape.width,
                        "original_left": shape.left,
                        "xx_yy_key": xx_yy_key
                    }
                    found_shapes_count += 1
                    
                # Added a more specific warning for shapes that match name but are group shapes (not common for bars)
                elif re.match(r"^bar_([A-Za-z]{1,2}\d+)$", shape.name) and shape.is_group:
                    logging.warning(f"  Shape '{shape.name}' on slide {slide_idx+1} matches naming convention but is a GroupShape. Skipping.")

    if found_shapes_count == 0:
        logging.warning("No shapes matching the 'bar_XXYY' naming convention were found in the presentation.")
    

def _update_percentage_shapes(all_tbls_flat: List[pd.DataFrame], lmap: Dict[str, int]):
    """
    Updates the width and color of identified percentage bar shapes based on Excel data.
    """
    GREEN = RGBColor(99, 195, 132)
    RED = RGBColor(255, 0, 0)
    WHITE = RGBColor(255, 255, 255) # For 0% or no change

    
    updated_shapes_count = 0
    for slide_id, shapes_on_slide in _percentage_shapes_metadata.items():
        for shape_name, shape_info in shapes_on_slide.items():
            shape_obj = shape_info["shape_obj"]
            original_width = shape_info["original_width"]
            original_left = shape_info["original_left"]
            xx_yy_key = shape_info["xx_yy_key"]

            
            try:
                # Extract prefix and number from XXYY key
                pat = re.compile(r"([A-Za-z]{1,2})(\d+)")
                match = pat.match(xx_yy_key)
                if not match:
                    logging.warning(f"    Invalid XXYY key format '{xx_yy_key}' for shape {shape_name}. Skipping shape update.")
                    continue

                prefix, number_str = match.groups()
                table_index = lmap.get(prefix)
                cell_index = int(number_str) - 1

                if table_index is not None and 0 <= table_index < len(all_tbls_flat):
                    df = all_tbls_flat[table_index]
                    df_ncols = df.shape[1]
                    df_nrows = df.shape[0]
                    df_row, df_col = divmod(cell_index, df_ncols)

                    if df_row < df_nrows and df_col < df_ncols:
                        raw_value = df.iat[df_row, df_col]
                        
                        
                        # Ensure raw_value is numeric for calculations
                        if not isinstance(raw_value, (int, float)):
                            logging.warning(f"    Non-numeric value '{raw_value}' for shape {shape_name}. Skipping shape update.")
                            continue

                        # Normalize percentage to 0.0-1.0 range (absolute value for width)
                        normalized_percentage_for_width = abs(raw_value)
                        
                        # Ensure normalized_percentage_for_width is between 0 and 1
                        normalized_percentage_for_width = max(0.0, min(1.0, normalized_percentage_for_width))

                        
                        # Set shape fill to solid
                        fill = shape_obj.fill
                        fill.solid()

                        # --- Apply new logic based on percentage value ---
                        if normalized_percentage_for_width >= 1.0: # 100% or over
                            shape_obj.width = original_width
                            shape_obj.left = original_left
                            fill.fore_color.rgb = GREEN
                            
                        elif normalized_percentage_for_width == 0.0:  # 0%
                            # Remove the shape from the slide if percentage is 0%
                            # `shape_obj` is the shape; we need to get its parent collection and call ._spTree.remove()
                            # The following is safe for python-pptx >= 0.6.17 (commonly used)
                            slide = shape_obj.part.slide
                            slide.shapes._spTree.remove(shape_obj._element)

                        else: # Between 0% and 100% (exclusive)
                            # CRITICAL FIX: Round the calculated width to an integer
                            new_width = round(original_width * normalized_percentage_for_width)
                            shape_obj.width = new_width
                            shape_obj.left = original_left # Maintain left edge for left-to-right growth
                            if raw_value > 0:
                                fill.fore_color.rgb = GREEN
                            else: # raw_value < 0
                                fill.fore_color.rgb = RED
                            
                        updated_shapes_count += 1

                    else:
                        logging.warning(f"    Cell index {df_row},{df_col} out of bounds for DataFrame for shape {shape_name}. Skipping shape update.")
                else:
                    logging.warning(f"    Table index {table_index} not found for shape {shape_name}. Skipping shape update.")
            except Exception as e:
                logging.error(f"    Error updating shape {shape_name}: {e}", exc_info=True)
    
    

def update_ppt_labels(pptx_path: str, output_path: str, table_data: Dict[str, List[pd.DataFrame]]):
    """
    Updates a PowerPoint presentation with data from Excel tables and date labels.
    Applies dynamic sizing and coloring to shapes named 'bar_XXYY'.
    """
    prs = Presentation(pptx_path)
    pat = re.compile(r"\{\{([A-Za-z]{1,2})(\d+)}}") # Pattern for {{XXYY}} placeholders
    
    # Custom order for mapping prefixes to tables (as per original script)
    custom_order = [
        "a", "aa", "b", "bb", "c", "cc", "d", "dd", "e", "ee", 
        "f", "g", "h", "i", "j", "k", "l", "ff", "m", "n", "o", "p", "q", 
        "gg", "r", "s", "t", "hh", "u", "v", "w", "ii", "x", "y", "z", "ab", 
        "ac", "ad", "ae", "af", "ag", "ah", "AA", "A", "AB", "BB", "B", "BC", "CC", "C", "CD", "DD", "D", "DE", "EE", "E", "EF", "FF", "F", "FG"
    ]

    # Flatten table_data into a single list for lmap indexing
    all_tbls_flat = []
    lmap = {} # Maps prefix (e.g., 'a') to its index in all_tbls_flat list
    for sheet_name in table_data:
        tables_in_sheet = table_data.get(sheet_name, [])
        for i, table_df in enumerate(tables_in_sheet):
            if len(all_tbls_flat) < len(custom_order):
                prefix = custom_order[len(all_tbls_flat)]
                lmap[prefix] = len(all_tbls_flat)
                all_tbls_flat.append(table_df)
            else:
                logging.warning(f"More tables found than prefixes in custom_order. Skipping table from sheet {sheet_name}.")


    # 1) Collect initial data for percentage shapes
    _collect_initial_shape_data(prs)

    # 2) Replace date tags first
    _replace_date_tags(prs)

    # 3) Replace numeric placeholders in text boxes and table cells
    for slide in prs.slides:
        for shp in iterate_all_shapes(slide.shapes):
            # Process text boxes
            if hasattr(shp, "text_frame"):
                for p in shp.text_frame.paragraphs:
                    for run in p.runs:
                        original_text = run.text
                        updated_text = original_text

                        for match in pat.findall(original_text):
                            prefix, number = match
                            table_index = lmap.get(prefix)
                            cell_index = int(number) - 1

                            if table_index is not None and 0 <= table_index < len(all_tbls_flat):
                                df = all_tbls_flat[table_index]
                                ncols = df.shape[1]
                                nrows = df.shape[0]
                                row, col = divmod(cell_index, ncols)

                                if row < nrows and col < ncols:
                                    value = df.iat[row, col]
                                    formatted_value, _ = _format_custom_value(prefix, value) 
                                    updated_text = updated_text.replace(f"{{{{{prefix}{number}}}}}", formatted_value)
                                    
                        if updated_text != original_text:
                            run.text = updated_text

            # Process table cells (only for text replacement, no background coloring here)
            if hasattr(shp, "has_table") and shp.has_table:
                for r_idx, row_obj in enumerate(shp.table.rows):
                    for c_idx, cell_obj in enumerate(row_obj.cells):
                        for p in cell_obj.text_frame.paragraphs:
                            for run in p.runs:
                                original_text = run.text
                                updated_text = original_text

                                for match in pat.findall(original_text):
                                    prefix, number = match
                                    table_index = lmap.get(prefix)
                                    df_cell_index = int(number) - 1

                                    if table_index is not None and 0 <= table_index < len(all_tbls_flat):
                                        df = all_tbls_flat[table_index]
                                        df_ncols = df.shape[1]
                                        df_nrows = df.shape[0]
                                        df_row, df_col = divmod(df_cell_index, df_ncols)

                                        if df_row < df_nrows and df_col < df_ncols:
                                            value = df.iat[df_row, df_col]
                                            formatted_value, _ = _format_custom_value(prefix, value)
                                            updated_text = updated_text.replace(f"{{{{{prefix}{number}}}}}", formatted_value)
                                            
                                if updated_text != original_text:
                                    run.text = updated_text
                                        
    # 4) Update percentage bar shapes
    _update_percentage_shapes(all_tbls_flat, lmap)

    prs.save(output_path)
    logging.info(f"Saved updated PowerPoint to {output_path}")